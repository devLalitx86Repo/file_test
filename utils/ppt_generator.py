from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE



class PPTGenerator:

    def __init__(self, template_file):
        self.prs = Presentation(template_file)


    def update_placeholder(self, slide, placeholder_name, value):
        # print(slide.shapes)
        for shape in slide.shapes:
            # print(shape.name)
            if shape.has_text_frame and shape.name.lower() == placeholder_name.lower():
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ''  # Clear text in existing runs
                # Add the new text to the first run of the first paragraph
                if text_frame.paragraphs:
                    paragraph = text_frame.paragraphs[0]
                    if paragraph.runs:
                        paragraph.runs[0].text = str(value)


    def get_chart(self, slide, chart_name):
        # print([shape.name for shape in slide.shapes])
        for shape in slide.shapes:
            if shape.has_chart and shape.name.lower() == chart_name.lower():
                chart = shape.chart
                return chart
        return None
        
    def update_chart(self, slide, chart_name, x_axis, y_axis, series_names=None):
        chart = self.get_chart(slide, chart_name)
        if chart is None:
            print(f"Chart '{chart_name}' not found in the slide.")
            return
        chart_data = CategoryChartData()
        chart_data.categories = x_axis
        # chart_data.add_series("", y_axis[0])
        # chart_data.add_series("", y_axis[1])
        for i, series_data in enumerate(y_axis):
            # series_name = "Series_"+str(i) 
            series_name = series_names[i] if series_names is not None and series_names[i] else f"Series_{i}"
            chart_data.add_series(series_name, series_data)
        # print(chart_data)
        # Save data in csv
        chart.replace_data(chart_data)
        # chart.chart_type = XL_CHART_TYPE.LINE



    def add_slide(self, layout):
        slide = self.prs.slides[layout]
        return slide


    def save_ppt(self, output_file):
        self.prs.save(output_file)



